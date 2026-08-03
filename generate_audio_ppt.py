"""
generate_audio_ppt.py
Reads a PowerPoint, extracts slide notes, synthesizes speech via AWS Polly,
and embeds the audio onto each slide for Storyline 360 import.

Run with no arguments to open the GUI. Pass CLI arguments for headless use.
"""

import argparse
import json
import logging
import os
import sys
import tempfile
import threading
import tkinter as tk
from pathlib import Path
from tkinter import filedialog, messagebox, ttk

import boto3
import windnd
from botocore.config import Config as BotoConfig
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
import urllib3

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

log = logging.getLogger("polly-tts")


# ── Tooltip helper ───────────────────────────────────────────────────────────
class ToolTip:
    """Tooltip that appears on hover over a widget."""

    def __init__(self, widget: tk.Widget, text: str):
        self.widget = widget
        self.text = text
        self.tip_window: tk.Toplevel | None = None
        widget.bind("<Enter>", self._show)
        widget.bind("<Leave>", self._hide)

    def _show(self, event=None):
        x = self.widget.winfo_rootx() + 20
        y = self.widget.winfo_rooty() + self.widget.winfo_height() + 4
        self.tip_window = tw = tk.Toplevel(self.widget)
        tw.wm_overrideredirect(True)
        tw.wm_geometry(f"+{x}+{y}")
        label = tk.Label(tw, text=self.text, background="#ffffe0",
                         relief="solid", borderwidth=1, font=("Segoe UI", 9))
        label.pack()

    def _hide(self, event=None):
        if self.tip_window:
            self.tip_window.destroy()
            self.tip_window = None

# ── Voice presets ────────────────────────────────────────────────────────────
VOICE_PRESETS = {
    "mandarin": {"voice": "Zhiyu", "lang": "cmn-CN"},
    "english": {"voice": "Matthew", "lang": "en-US"},
    "korean": {"voice": "Seoyeon", "lang": "ko-KR"},
}

PRESET_LABELS = {
    "mandarin": "Mandarin (Zhiyu)",
    "english": "English (Matthew)",
    "korean": "Korean (Seoyeon)",
}

CONFIG_FILE = "config.json"


# ── Config helpers ───────────────────────────────────────────────────────────
def get_config_path() -> Path:
    """Return the path to config.json next to the executable or script."""
    if getattr(sys, "frozen", False):
        base = Path(sys.executable).parent
    else:
        base = Path(__file__).parent
    return base / CONFIG_FILE


def load_config() -> dict:
    """Load config.json if it exists."""
    cfg_path = get_config_path()
    if cfg_path.exists():
        with open(cfg_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_config(config: dict) -> None:
    """Write config back to config.json."""
    cfg_path = get_config_path()
    with open(cfg_path, "w", encoding="utf-8") as f:
        json.dump(config, f, indent=4, ensure_ascii=False)


_EMPTY_PRONUNCIATIONS = {k: [] for k in VOICE_PRESETS}


def load_pronunciations(config: dict | None = None) -> dict:
    """Return per-voice pronunciation correction lists from config.

    Returns a deep copy so callers can mutate safely without affecting the
    module-level default or the on-disk config.
    """
    if config is None:
        config = load_config()
    src = config.get("pronunciations", _EMPTY_PRONUNCIATIONS)
    return {k: [dict(item) for item in v] for k, v in src.items()}


def save_pronunciations(pronunciations: dict, config: dict | None = None) -> None:
    """Persist pronunciation corrections into config.json."""
    if config is None:
        config = load_config()
    config["pronunciations"] = pronunciations
    save_config(config)


def create_polly_client(config: dict = None):
    """Create a boto3 Polly client, optionally using credentials from config."""
    kwargs = {
        "config": BotoConfig(retries={"max_attempts": 3, "mode": "adaptive"}),
    }
    if config:
        aws = config.get("aws", {})
        key_id = aws.get("access_key_id", "")
        secret = aws.get("secret_access_key", "")
        if key_id and secret and not key_id.startswith("YOUR_"):
            kwargs["aws_access_key_id"] = key_id
            kwargs["aws_secret_access_key"] = secret
            if aws.get("region"):
                kwargs["region_name"] = aws["region"]
        if aws.get("verify_ssl", False) is False:
            kwargs["verify"] = False
    return boto3.client("polly", **kwargs)


# ── Core logic ───────────────────────────────────────────────────────────────
def parse_slide_range(spec: str, total: int) -> list[int]:
    """Parse '1,3,5-8' into sorted 0-indexed slide indices."""
    indices: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            lo, hi = part.split("-", 1)
            lo, hi = int(lo), int(hi)
            if lo < 1 or hi > total or lo > hi:
                raise ValueError(f"Invalid range {part} for presentation with {total} slides")
            indices.update(range(lo - 1, hi))
        else:
            n = int(part)
            if n < 1 or n > total:
                raise ValueError(f"Slide {n} out of range (1-{total})")
            indices.add(n - 1)
    return sorted(indices)


def _is_hangul(ch: str) -> bool:
    """Return True if *ch* is a Korean Hangul syllable or jamo."""
    cp = ord(ch)
    return (
        0xAC00 <= cp <= 0xD7AF   # Hangul Syllables
        or 0x1100 <= cp <= 0x11FF  # Hangul Jamo
        or 0x3130 <= cp <= 0x318F  # Hangul Compatibility Jamo
    )


def _is_cjk_ideograph(ch: str) -> bool:
    """Return True if *ch* is a CJK Unified Ideograph (Chinese/Japanese)."""
    cp = ord(ch)
    return 0x4E00 <= cp <= 0x9FFF


def detect_language(filename: str) -> str:
    """Heuristically detect voice preset from the filename.

    Returns ``"mandarin"``, ``"korean"``, or ``"english"`` based on the
    character composition of the filename stem.
    """
    text = Path(filename).stem
    hangul = cjk = total = 0
    for ch in text:
        if ch.isspace():
            continue
        total += 1
        if _is_hangul(ch):
            hangul += 1
        elif _is_cjk_ideograph(ch):
            cjk += 1
    if total == 0:
        return "mandarin"
    hangul_ratio = hangul / total
    cjk_ratio = cjk / total
    if hangul_ratio > 0.3:
        return "korean"
    if cjk_ratio > 0.3:
        return "mandarin"
    return "english"


def apply_corrections(text: str, corrections: list[dict]) -> str:
    """Apply pronunciation corrections to *text* using longest-match-first.

    *corrections* is a list of dicts with keys ``original`` and ``replacement``.
    Overlapping matches are resolved by preferring the longest ``original`` string.
    """
    if not corrections:
        return text
    # Sort longest-original-first so greediest match wins
    active = sorted(corrections, key=lambda c: len(c["original"]), reverse=True)
    for entry in active:
        text = text.replace(entry["original"], entry["replacement"])
    return text


def synthesize_audio(polly_client, text: str, output_path: str,
                     voice_id: str, engine: str, lang_code: str) -> bool:
    """Call AWS Polly and save the result as MP3. Returns True on success."""
    response = polly_client.synthesize_speech(
        Engine=engine,
        OutputFormat="mp3",
        Text=text,
        LanguageCode=lang_code,
        VoiceId=voice_id,
    )
    if "AudioStream" in response:
        with open(output_path, "wb") as f:
            f.write(response["AudioStream"].read())
        return True
    return False


def embed_audio_on_slide(slide, audio_path: str) -> None:
    """Embed an MP3 onto a slide using add_movie() then patch timing for autoplay."""
    from lxml import etree

    slide.shapes.add_movie(
        audio_path,
        left=Inches(-2.0),
        top=Inches(-2.0),
        width=Inches(1.0),
        height=Inches(1.0),
        poster_frame_image=None,
        mime_type="audio/mpeg",
    )

    timing = slide._element.find(qn("p:timing"))
    if timing is not None:
        for cond in timing.iter(qn("p:cond")):
            if cond.get("delay") == "indefinite":
                cond.set("delay", "0")


def normalize_audio_for_storyline(output_path: Path, callback=None) -> bool:
    """Open in PowerPoint via COM, nudge each shape to trigger media normalization.

    This step is required for Storyline 360 to properly play embedded audio.
    PowerPoint normalizes the media format when the file is opened and shapes
    are interacted with.
    """
    abs_path = str(output_path.resolve())
    log.info("Normalizing audio for Storyline via PowerPoint COM ...")
    if callback:
        callback("Normalizing audio for Storyline compatibility...")
    try:
        import time

        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        if callback:
            callback("Opening file in PowerPoint...")
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        pres = ppt.Presentations.Open(abs_path, ReadOnly=False)

        total_slides = pres.Slides.Count
        if callback:
            callback(f"Normalizing audio on slide 1 of {total_slides} (temp: first slide only)...")

        # TEMP: only normalize first slide for faster testing
        for slide_idx in range(1, 2):
            slide = pres.Slides(slide_idx)
            slide.Select()
            has_media = False
            for shape_idx in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(shape_idx)
                try:
                    mt = shape.MediaType
                except Exception:
                    mt = 0
                if mt == 0:
                    continue
                has_media = True
                orig_left = shape.Left
                shape.Left = orig_left + 1
                time.sleep(0.05)
                shape.Left = orig_left
                log.debug("Nudged shape '%s' on slide %d", shape.Name, slide_idx)
            if callback:
                callback(f"Slide {slide_idx}/{total_slides}: {'audio normalized' if has_media else 'no audio'}")

        pres.Save()
        pres.Close()
        ppt.Quit()
        pythoncom.CoUninitialize()
        log.info("PowerPoint normalization complete")
        if callback:
            callback("Audio normalization complete — ready for Storyline import")
        return True
    except Exception as e:
        err = str(e)
        log.warning("PowerPoint COM normalization failed: %s", err)
        if callback:
            if "not installed" in err.lower() or "cannot be found" in err.lower() or "class not registered" in err.lower() or "429" in err:
                callback("WARNING: PowerPoint not installed — audio embedded but may not work in Storyline. Install PowerPoint to enable normalization.")
            else:
                callback(f"WARNING: Normalization failed: {err}")
        try:
            pres.Close()
        except Exception:
            pass
        try:
            ppt.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
        return False


def process_pptx(input_path: Path, output_path: Path, voice_id: str, engine: str,
                 lang_code: str, slides_spec: str | None = None,
                 callback=None, config: dict = None,
                 pronunciations: list[dict] | None = None) -> dict:
    """
    Main processing pipeline. Returns a summary dict.
    callback(message) is called with status updates.
    *pronunciations* is the correction list for the active voice; if provided
    a *copy* of each slide's notes is corrected before being sent to Polly.
    """
    def emit(msg):
        log.info(msg)
        if callback:
            callback(msg)

    prs = Presentation(str(input_path))
    total = len(prs.slides)
    emit(f"Loaded '{input_path.name}' — {total} slide(s)")

    if slides_spec:
        slide_indices = parse_slide_range(slides_spec, total)
    else:
        slide_indices = list(range(total))

    emit(f"Processing {len(slide_indices)} slide(s): "
         + ", ".join(str(i + 1) for i in slide_indices))

    config = config or load_config()
    polly_client = create_polly_client(config)
    emit("AWS Polly client ready")

    success_count = skip_count = fail_count = 0

    with tempfile.TemporaryDirectory(prefix="polly_tts_") as tmp_dir:
        for i, idx in enumerate(slide_indices):
            slide = prs.slides[idx]
            slide_num = idx + 1

            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            if not notes_text:
                emit(f"Slide {slide_num}: no notes — skipping")
                skip_count += 1
                continue

            tts_text = apply_corrections(notes_text, pronunciations or [])

            emit(f"Slide {slide_num}: synthesizing {len(tts_text)} chars ... ({i + 1}/{len(slide_indices)})")
            audio_path = os.path.join(tmp_dir, f"slide_{slide_num}.mp3")

            try:
                ok = synthesize_audio(polly_client, tts_text, audio_path,
                                      voice_id, engine, lang_code)
                if not ok:
                    emit(f"Slide {slide_num}: Polly returned no audio stream")
                    fail_count += 1
                    continue
            except Exception as e:
                emit(f"Slide {slide_num}: Polly API error: {e}")
                fail_count += 1
                continue

            embed_audio_on_slide(slide, audio_path)
            emit(f"Slide {slide_num}: audio embedded")
            success_count += 1

    prs.save(str(output_path))
    emit(f"Saved: {output_path}")

    normalize_audio_for_storyline(output_path, callback=callback)

    emit(f"Done — {success_count} slide(s) narrated, {skip_count} skipped (no notes), {fail_count} failed")
    return {"success": success_count, "skipped": skip_count, "failed": fail_count}


# ── CLI ──────────────────────────────────────────────────────────────────────
def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Add AWS Polly TTS narration to PowerPoint slides.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "voice presets:\n"
            "  mandarin   -> Zhiyu     / cmn-CN\n"
            "  english    -> Matthew   / en-US\n"
            "  korean     -> Seoyeon   / ko-KR\n"
            "\nRun with no arguments to open the GUI."
        ),
    )
    p.add_argument("input", nargs="?", help="Path to the input .pptx file")
    p.add_argument("-o", "--output", help="Output file path (default: narrated_<input>)")
    p.add_argument("--voice", default="Zhiyu", help="Polly voice ID (default: Zhiyu)")
    p.add_argument("--engine", choices=["standard", "neural"], default="neural",
                   help="TTS engine (default: neural)")
    p.add_argument("--lang", default="cmn-CN", help="Language code (default: cmn-CN)")
    p.add_argument("--preset", choices=list(VOICE_PRESETS), default=None,
                   help="Use a named voice preset instead of --voice/--lang")
    p.add_argument("--slides", default=None,
                   help="Slide numbers or ranges, e.g. '1,3,5-8' (default: all)")
    p.add_argument("--dry-run", action="store_true",
                   help="Show extracted notes without calling AWS")
    p.add_argument("--keep-audio", action="store_true",
                   help="Keep temp MP3 files after processing")
    p.add_argument("--gui", action="store_true", help="Force open the GUI")
    p.add_argument("-v", "--verbose", action="store_true", help="Enable debug logging")
    return p


def cli_main(args=None):
    parsed = build_parser().parse_args(args)

    level = logging.DEBUG if parsed.verbose else logging.INFO
    logging.basicConfig(level=level, format="[%(levelname)s] %(message)s")

    if not parsed.input and not parsed.gui:
        gui_main()
        return

    voice_id = parsed.voice
    lang_code = parsed.lang
    if parsed.preset:
        preset = VOICE_PRESETS[parsed.preset]
        voice_id = preset["voice"]
        lang_code = preset["lang"]

    input_path = Path(parsed.input)
    if not input_path.exists():
        log.error("File not found: %s", input_path)
        sys.exit(1)

    if parsed.output:
        output_path = Path(parsed.output)
    else:
        stem = input_path.stem
        suffix = input_path.suffix
        output_path = input_path.with_name(f"narrated_{stem}{suffix}")
        i = 1
        while output_path.exists():
            output_path = input_path.with_name(f"narrated_{stem}_{i}{suffix}")
            i += 1

    if parsed.dry_run:
        prs = Presentation(str(input_path))
        total = len(prs.slides)
        indices = parse_slide_range(parsed.slides, total) if parsed.slides else list(range(total))
        for idx in indices:
            slide = prs.slides[idx]
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                preview = notes[:200].replace("\n", " ")
                log.info("Slide %d — %d chars: %s%s", idx + 1, len(notes),
                         preview, "..." if len(notes) > 200 else "")
            else:
                log.info("Slide %d — (no notes)", idx + 1)
        log.info("Dry run complete. No audio generated.")
        return

    process_pptx(input_path, output_path, voice_id, parsed.engine, lang_code,
                 slides_spec=parsed.slides)


# ── Pronunciation Corrections Dialog ──────────────────────────────────────────
class CorrectionsDialog:
    """Separate window for managing per-voice pronunciation corrections."""

    def __init__(self, parent: tk.Tk, config: dict, initial_voice: str = "mandarin"):
        self.config = config
        self.pronunciations = load_pronunciations(config)
        self.active_voice = tk.StringVar(value=initial_voice)

        self.win = tk.Toplevel(parent)
        self.win.title("Pronunciation Corrections")
        self.win.geometry("560x480")
        self.win.resizable(False, False)
        self.win.transient(parent)
        self.win.grab_set()

        self._rows: list[dict] = []
        self._build_ui()
        self._refresh_rows()

        self.active_voice.trace_add("write", lambda *_: self._refresh_rows())

    # ── UI construction ──────────────────────────────────────────────────
    def _build_ui(self):
        pad = {"padx": 8, "pady": 4}

        # Voice selector
        top = ttk.Frame(self.win)
        top.pack(fill="x", **pad)
        ttk.Label(top, text="Language:").pack(side="left", padx=(0, 6))
        voices = [f"{k} — {PRESET_LABELS[k]}" for k in VOICE_PRESETS]
        cb = ttk.Combobox(top, textvariable=self.active_voice, state="readonly",
                          width=30, values=voices)
        cb.pack(side="left")
        # Set initial selection to match the main window's language
        try:
            initial_idx = list(VOICE_PRESETS.keys()).index(self.active_voice.get())
            cb.current(initial_idx)
        except ValueError:
            cb.current(0)

        # Scrollable list area
        list_frame = ttk.LabelFrame(self.win, text="Corrections", padding=4)
        list_frame.pack(fill="both", expand=True, **pad)

        self._canvas = tk.Canvas(list_frame, highlightthickness=0)
        scrollbar = ttk.Scrollbar(list_frame, orient="vertical",
                                  command=self._canvas.yview)
        self._inner = ttk.Frame(self._canvas)

        self._inner.bind("<Configure>",
                         lambda _: self._canvas.configure(scrollregion=self._canvas.bbox("all")))
        self._canvas.create_window((0, 0), window=self._inner, anchor="nw")
        self._canvas.configure(yscrollcommand=scrollbar.set)

        scrollbar.pack(side="right", fill="y")
        self._canvas.pack(side="left", fill="both", expand=True)

        # Column headers
        hdr = ttk.Frame(self._inner)
        hdr.pack(fill="x", pady=(0, 2))
        ttk.Label(hdr, text="Find", width=18, anchor="w").pack(side="left", padx=(4, 2))
        ttk.Label(hdr, text="", width=3).pack(side="left")
        ttk.Label(hdr, text="Replace With", width=18, anchor="w").pack(side="left", padx=(2, 2))

        # Add row
        add_frame = ttk.Frame(self.win)
        add_frame.pack(fill="x", **pad)

        # Example hint
        hint_frame = ttk.Frame(self.win)
        hint_frame.pack(fill="x", padx=8, pady=(0, 4))
        ttk.Label(hint_frame, text="e.g.  Find: 你好   Replace with: ni hao",
                  font=("Segoe UI", 8, "italic"), foreground="#666666").pack(anchor="w")

        self._add_orig = tk.StringVar()
        self._add_repl = tk.StringVar()
        ttk.Label(add_frame, text="Add:").pack(side="left", padx=(0, 4))
        ttk.Entry(add_frame, textvariable=self._add_orig, width=16).pack(
            side="left", padx=(0, 4))
        ttk.Label(add_frame, text="→").pack(side="left", padx=(0, 4))
        ttk.Entry(add_frame, textvariable=self._add_repl, width=16).pack(
            side="left", padx=(0, 8))
        ttk.Button(add_frame, text="Add", command=self._on_add).pack(side="left")

        # Bottom buttons
        btn_frame = ttk.Frame(self.win)
        btn_frame.pack(fill="x", **pad)
        ttk.Button(btn_frame, text="Preview Changes", command=self._on_preview).pack(side="left")
        ttk.Button(btn_frame, text="Save & Close", command=self._on_save_close).pack(
            side="right")

    # ── Row management ───────────────────────────────────────────────────
    def _current_list(self) -> list[dict]:
        key = self.active_voice.get().split(" — ")[0]
        return self.pronunciations.setdefault(key, [])

    def _refresh_rows(self):
        for entry in self._rows:
            entry["frame"].destroy()
        self._rows.clear()

        corrections = self._current_list()
        if corrections:
            for i, item in enumerate(corrections):
                self._add_row(i, item)
        else:
            self._show_empty_state()
        self._canvas.yview_moveto(0)

    def _show_empty_state(self):
        frame = ttk.Frame(self._inner)
        frame.pack(fill="x", pady=(20, 10))
        ttk.Label(frame, text="No pronunciation rules yet.",
                  font=("Segoe UI", 9, "italic")).pack()
        ttk.Label(frame, text="Add a rule below to correct how words are spoken.",
                  font=("Segoe UI", 9)).pack()

    def _add_row(self, idx: int, item: dict):
        frame = ttk.Frame(self._inner)
        frame.pack(fill="x", pady=1)

        orig_var = tk.StringVar(value=item["original"])
        repl_var = tk.StringVar(value=item["replacement"])

        orig_entry = ttk.Entry(frame, textvariable=orig_var, width=18)
        orig_entry.pack(side="left", padx=(4, 2))
        ttk.Label(frame, text="→").pack(side="left", padx=(0, 2))
        repl_entry = ttk.Entry(frame, textvariable=repl_var, width=18)
        repl_entry.pack(side="left", padx=(2, 2))

        del_btn = ttk.Button(frame, text="✕", width=3,
                             command=lambda: self._delete(idx))
        del_btn.pack(side="left", padx=(8, 4))

        entry = {
            "frame": frame,
            "orig_var": orig_var,
            "repl_var": repl_var,
            "orig_entry": orig_entry,
            "repl_entry": repl_entry,
        }
        self._rows.append(entry)

        orig_var.trace_add("write", lambda *_a, i=idx, v=orig_var: self._sync_item(i, "original", v.get()))
        repl_var.trace_add("write", lambda *_a, i=idx, v=repl_var: self._sync_item(i, "replacement", v.get()))

    def _sync_item(self, idx: int, field: str, value: str):
        lst = self._current_list()
        if idx < len(lst):
            lst[idx][field] = value

    def _delete(self, idx: int):
        lst = self._current_list()
        if idx < len(lst):
            del lst[idx]
        self._refresh_rows()

    def _on_add(self):
        orig = self._add_orig.get().strip()
        repl = self._add_repl.get().strip()
        if not orig or not repl:
            messagebox.showwarning("Missing values", "Both fields are required.",
                                   parent=self.win)
            return
        self._current_list().append({
            "original": orig,
            "replacement": repl,
        })
        self._add_orig.set("")
        self._add_repl.set("")
        self._refresh_rows()

    # ── Preview ──────────────────────────────────────────────────────────
    def _on_preview(self):
        input_file = self._get_input_file()
        if not input_file:
            return

        key = self.active_voice.get().split(" — ")[0]
        display_name = PRESET_LABELS.get(key, key)
        corrections = self.pronunciations.get(key, [])

        try:
            prs = Presentation(str(input_file))
        except Exception as e:
            messagebox.showerror("Error", f"Cannot open PPTX:\n{e}", parent=self.win)
            return

        preview_win = tk.Toplevel(self.win)
        preview_win.title(f"Preview — {display_name}")
        preview_win.geometry("640x400")
        preview_win.transient(self.win)

        txt = tk.Text(preview_win, wrap="word", font=("Consolas", 9), state="normal")
        txt.pack(fill="both", expand=True, padx=8, pady=8)
        txt.tag_configure("highlight", background="#ffff80", foreground="#000000")
        txt.tag_configure("slide_hdr", font=("Consolas", 9, "bold"))

        total_replacements = 0
        PREVIEW_LEN = 60

        for idx, slide in enumerate(prs.slides):
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            if not notes:
                continue

            corrected = apply_corrections(notes, corrections)
            if corrected == notes:
                continue

            # Count replacements
            active = sorted(corrections, key=lambda c: len(c["original"]), reverse=True)
            slide_count = 0
            tmp = notes
            for entry in active:
                before = tmp
                tmp = tmp.replace(entry["original"], entry["replacement"])
                slide_count += before.count(entry["original"])
            total_replacements += slide_count

            hdr = f"── Slide {idx + 1} ({slide_count} replacement{'s' if slide_count != 1 else ''}) ──\n"
            txt.insert("end", hdr, "slide_hdr")

            # Before line
            before_preview = notes[:PREVIEW_LEN]
            if len(notes) > PREVIEW_LEN:
                before_preview += "..."
            txt.insert("end", f"  Before: {before_preview}\n")

            # After line with highlighted replacements
            txt.insert("end", "  After:  ")
            corr_trunc = corrected[:PREVIEW_LEN]
            spans = self._compute_highlight_spans(notes, corrected, active)
            pos = 0
            for span_start, span_end in spans:
                if span_start > pos:
                    txt.insert("end", corr_trunc[pos:span_start])
                txt.insert("end", corr_trunc[span_start:span_end], "highlight")
                pos = span_end
            if pos < len(corr_trunc):
                txt.insert("end", corr_trunc[pos:])
            if len(corrected) > PREVIEW_LEN:
                txt.insert("end", "...")
            txt.insert("end", "\n\n")

        if total_replacements == 0:
            txt.insert("end", "No replacements to preview — no notes match any correction.\n")

        preview_win.title(f"Preview — {display_name} ({total_replacements} replacement"
                          f"{'s' if total_replacements != 1 else ''})")
        txt.configure(state="disabled")

    @staticmethod
    def _compute_highlight_spans(original: str, corrected: str,
                                 corrections: list[dict]) -> list[tuple[int, int]]:
        """Return (start, end) spans in *corrected* that were changed by replacements.

        Walks a working copy of *original*, applying each replacement
        sequentially.  Positions found via ``str.find`` on the mutated copy
        are already in corrected-text coordinates, so no offset tracking
        is needed.
        """
        spans: list[tuple[int, int]] = []
        text = original
        for entry in corrections:
            orig_str = entry["original"]
            repl_str = entry["replacement"]
            i = 0
            while i < len(text):
                idx = text.find(orig_str, i)
                if idx == -1:
                    break
                spans.append((idx, idx + len(repl_str)))
                text = text[:idx] + repl_str + text[idx + len(orig_str):]
                i = idx + len(repl_str)
        spans.sort()
        return spans

    def _get_input_file(self) -> str | None:
        app = getattr(self.win.master, "_ppttts_app", None)
        if app:
            val = app.input_path.get().strip()
            if val and Path(val).exists():
                return val
        messagebox.showwarning("No file", "Select a PowerPoint file first.",
                               parent=self.win)
        return None

    # ── Save & close ─────────────────────────────────────────────────────
    def _on_save_close(self):
        save_pronunciations(self.pronunciations, self.config)
        self.win.destroy()


# ── GUI ──────────────────────────────────────────────────────────────────────
class PPTTTSApp:
    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("PowerPoint Narrator")
        self.root.geometry("640x580")
        self.root.resizable(False, False)

        self.input_path = tk.StringVar()
        self.preset_var = tk.StringVar(value="mandarin")
        self.slides_var = tk.StringVar()
        self.processing = False
        self.cancel_flag = threading.Event()
        self.config = load_config()
        self.root._ppttts_app = self  # let the dialog find us via root

        self._build_ui()

    def _build_ui(self):
        pad = {"padx": 12, "pady": 4}

        # ── File selection ───────────────────────────────────────────────
        file_frame = ttk.LabelFrame(self.root, text="PowerPoint File", padding=8)
        file_frame.pack(fill="x", **pad)

        self.file_entry = ttk.Entry(file_frame, textvariable=self.input_path, width=60)
        self.file_entry.pack(side="left", fill="x", expand=True, padx=(0, 8))
        self.browse_btn = ttk.Button(file_frame, text="Browse...", command=self._browse_file)
        self.browse_btn.pack(side="right")

        # Drag-and-drop on the entire window
        windnd.hook_dropfiles(self.root, func=self._on_drop)
        ToolTip(self.file_entry, "Type a path or drag a .pptx file anywhere on this window")
        ToolTip(self.browse_btn, "Select a PowerPoint file (.pptx)")

        # ── Options ─────────────────────────────────────────────────────
        opt_frame = ttk.LabelFrame(self.root, text="Options", padding=8)
        opt_frame.pack(fill="x", **pad)

        ttk.Label(opt_frame, text="Language:").grid(row=0, column=0, sticky="w", padx=(0, 8))
        preset_combo = ttk.Combobox(
            opt_frame, textvariable=self.preset_var, state="readonly", width=28,
            values=[f"{k} — {v}" for k, v in PRESET_LABELS.items()],
        )
        preset_combo.grid(row=0, column=1, sticky="w")
        self._preset_keys = list(PRESET_LABELS.keys())
        self._preset_combo = preset_combo
        preset_combo.current(0)
        ToolTip(preset_combo, "Voice and language for narration (auto-detected from filename)")

        ttk.Label(opt_frame, text="Slides:").grid(
            row=1, column=0, sticky="w", padx=(0, 8), pady=(6, 0))
        slides_entry = ttk.Entry(opt_frame, textvariable=self.slides_var, width=28)
        slides_entry.grid(row=1, column=1, sticky="w", pady=(6, 0))
        ttk.Label(opt_frame, text="optional, e.g. 1,3,5-8").grid(
            row=1, column=2, sticky="w", padx=(6, 0), pady=(6, 0))
        ToolTip(slides_entry, "Leave blank for all slides, or specify e.g. 1,3,5-8")

        self.pronunciation_btn = ttk.Button(opt_frame, text="Pronunciation...",
                   command=self._open_corrections)
        self.pronunciation_btn.grid(
            row=2, column=0, columnspan=3, sticky="w", pady=(8, 0))
        ToolTip(self.pronunciation_btn, "Define rules to correct how words are spoken (e.g. acronyms, names)")

        # ── Generate button ─────────────────────────────────────────────
        btn_frame = ttk.Frame(self.root)
        btn_frame.pack(fill="x", padx=12, pady=8)

        self.generate_btn = tk.Button(
            btn_frame, text="Generate Audio", command=self._on_generate,
            font=("Segoe UI", 11, "bold"), bg="#4CAF50", fg="white",
            activebackground="#45a049", activeforeground="white",
            relief="raised", padx=20, pady=6, cursor="hand2")
        self.generate_btn.pack(side="left")
        ToolTip(self.generate_btn, "Generate narration audio for all slides")

        self.cancel_btn = tk.Button(
            btn_frame, text="Cancel", command=self._on_cancel,
            font=("Segoe UI", 10), bg="#f44336", fg="white",
            activebackground="#d32f2f", activeforeground="white",
            relief="raised", padx=12, pady=6, cursor="hand2")
        self.cancel_btn.pack(side="left", padx=(8, 0))
        self.cancel_btn.pack_forget()

        # ── Progress bar ───────────────────────────────────────────────
        self.progress_var = tk.DoubleVar(value=0)
        self.progress_frame = ttk.Frame(self.root)
        self.progress_frame.pack(fill="x", padx=12, pady=(0, 4))

        self.progress_label = ttk.Label(self.progress_frame, text="")
        self.progress_label.pack(side="left", padx=(0, 8))

        self.progress_bar = ttk.Progressbar(
            self.progress_frame, variable=self.progress_var,
            maximum=100, mode="determinate")
        self.progress_bar.pack(side="left", fill="x", expand=True)
        self.progress_frame.pack_forget()

        # ── Log area ────────────────────────────────────────────────────
        log_frame = ttk.LabelFrame(self.root, text="Status", padding=8)
        log_frame.pack(fill="both", expand=True, **pad)

        self.log_text = tk.Text(log_frame, height=12, state="disabled",
                                wrap="word", font=("Consolas", 9))
        scrollbar = ttk.Scrollbar(log_frame, orient="vertical",
                                  command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side="right", fill="y")
        self.log_text.pack(fill="both", expand=True)

    def _browse_file(self):
        path = filedialog.askopenfilename(
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")])
        if path:
            self._set_input_file(path)

    def _on_drop(self, file_paths):
        path = file_paths[0].decode("utf-8") if isinstance(file_paths[0], bytes) else file_paths[0]
        if path.lower().endswith(".pptx"):
            self._set_input_file(path)
        else:
            messagebox.showwarning("Invalid file", "Please drop a .pptx file.")

    def _set_input_file(self, path: str):
        self.input_path.set(path)
        detected = detect_language(path)
        idx = self._preset_keys.index(detected) if detected in self._preset_keys else 0
        self._preset_combo.current(idx)

    def _resolve_preset_key(self) -> str:
        combo_val = self.preset_var.get()
        for k in self._preset_keys:
            if combo_val.startswith(k):
                return k
        return self._preset_keys[0]

    def _open_corrections(self):
        CorrectionsDialog(self.root, self.config, initial_voice=self._resolve_preset_key())

    def _log(self, msg: str):
        if not hasattr(self, "_log_buffer"):
            self._log_buffer = []
        self._log_buffer.append(msg)
        full = "\n".join(self._log_buffer)
        self.root.after(0, _append_set, self, full)

    def _on_generate(self):
        if self.processing:
            return

        input_file = self.input_path.get().strip()
        if not input_file:
            messagebox.showwarning("Missing file", "Please select a PowerPoint file.")
            return

        if not Path(input_file).exists():
            messagebox.showerror("File not found", f"Cannot find:\n{input_file}")
            return

        preset = VOICE_PRESETS[self._resolve_preset_key()]

        input_path = Path(input_file)
        output_path = input_path.with_name(f"narrated_{input_path.stem}{input_path.suffix}")
        i = 1
        while output_path.exists():
            output_path = input_path.with_name(f"narrated_{input_path.stem}_{i}{input_path.suffix}")
            i += 1

        slides_spec = self.slides_var.get().strip() or None

        self.processing = True
        self.cancel_flag.clear()
        self.generate_btn.configure(state="disabled")
        self.cancel_btn.pack(side="left", padx=(8, 0))
        self.browse_btn.configure(state="disabled")
        self._preset_combo.configure(state="disabled")
        self.pronunciation_btn.configure(state="disabled")
        self.progress_var.set(0)
        self.progress_label.configure(text="Starting...")
        self.progress_frame.pack(fill="x", padx=12, pady=(0, 4))
        self._log_buffer = []

        def worker():
            preset_key = self._resolve_preset_key()
            try:
                if slides_spec:
                    prs = Presentation(str(input_path))
                    parse_slide_range(slides_spec, len(prs.slides))

                all_pronunciations = load_pronunciations(self.config)
                voice_corrections = all_pronunciations.get(preset_key, [])

                def progress_callback(msg):
                    self._log(msg)
                    if self.cancel_flag.is_set():
                        raise InterruptedError("Cancelled by user")
                    if "synthesizing" in msg.lower() or "processing" in msg.lower():
                        try:
                            current = msg.split("(")[1].split(")")[0]
                            parts = current.split("/")
                            done = int(parts[0])
                            total = int(parts[1])
                            pct = (done / total) * 100
                            self.root.after(0, lambda p=pct: self.progress_var.set(p))
                            self.root.after(0, lambda d=done, t=total: self.progress_label.configure(text=f"Processing slide {d} of {t}"))
                        except Exception:
                            pass

                process_pptx(
                    input_path, output_path,
                    voice_id=preset["voice"],
                    engine="neural",
                    lang_code=preset["lang"],
                    slides_spec=slides_spec,
                    callback=progress_callback,
                    pronunciations=voice_corrections,
                )
                self.root.after(0, lambda: self.progress_var.set(100))
                self.root.after(0, lambda: self.progress_label.configure(text="Complete"))
                self.root.after(1000, lambda: self.progress_frame.pack_forget())
                self.root.after(0, lambda: self._show_completion_dialog(output_path))
            except InterruptedError:
                self._log("Cancelled by user")
                self.root.after(0, lambda: self.progress_label.configure(text="Cancelled"))
                self.root.after(1000, lambda: self.progress_frame.pack_forget())
            except Exception as e:
                self._log(f"ERROR: {e}")
                self.root.after(0, lambda: self.progress_frame.pack_forget())
                self.root.after(0, lambda: messagebox.showerror(
                    "Error", str(e)))
            finally:
                self.processing = False
                self.cancel_flag.clear()
                self.root.after(0, lambda: self.generate_btn.configure(state="normal"))
                self.root.after(0, lambda: self.cancel_btn.pack_forget())
                self.root.after(0, lambda: self.cancel_btn.configure(state="normal"))
                self.root.after(0, lambda: self.browse_btn.configure(state="normal"))
                self.root.after(0, lambda: self._preset_combo.configure(state="readonly"))
                self.root.after(0, lambda: self.pronunciation_btn.configure(state="normal"))

        threading.Thread(target=worker, daemon=True).start()

    def _on_cancel(self):
        if self.processing:
            self.cancel_flag.set()
            self.cancel_btn.configure(state="disabled")
            self.progress_label.configure(text="Stopping...")

    def _show_completion_dialog(self, output_path: Path):
        dialog = tk.Toplevel(self.root)
        dialog.title("Complete")
        dialog.geometry("360x180")
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()

        pad = {"padx": 16, "pady": 6}
        ttk.Label(dialog, text="Audio generation complete!",
                  font=("Segoe UI", 11, "bold")).pack(**pad)
        ttk.Label(dialog, text=f"Saved to:\n{output_path}",
                  font=("Segoe UI", 9), wraplength=320).pack(padx=16, pady=(0, 12))

        btn_frame = ttk.Frame(dialog)
        btn_frame.pack(fill="x", padx=16, pady=(0, 12))
        ttk.Button(btn_frame, text="Open File",
                   command=lambda: [os.startfile(str(output_path)), dialog.destroy()]).pack(
            side="left", padx=(0, 8))
        ttk.Button(btn_frame, text="Open Folder",
                   command=lambda: [os.startfile(str(output_path.parent)), dialog.destroy()]).pack(
            side="left", padx=(0, 8))
        ttk.Button(btn_frame, text="Close", command=dialog.destroy).pack(side="left")


def _append_set(app, full_text):
    app.log_text.configure(state="normal")
    app.log_text.delete("1.0", "end")
    app.log_text.insert("end", full_text)
    app.log_text.configure(state="disabled")
    app.log_text.see("end")


def gui_main():
    root = tk.Tk()
    PPTTTSApp(root)
    root.mainloop()


if __name__ == "__main__":
    cli_main()
